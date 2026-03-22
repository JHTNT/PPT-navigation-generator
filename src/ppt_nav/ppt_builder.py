from __future__ import annotations

import unicodedata
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation as PresentationFactory
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from ppt_nav.outline import Outline, OutlineItem, SlidePlanEntry


class PresentationBuilder:
    def __init__(self, font_size: Optional[float] = None) -> None:
        self.nav_side_margin = Inches(0)
        self.nav_top_margin = Inches(0)
        self.body_margin_top = Inches(0.3)
        self.body_side_margin = Inches(0.5)

        # Theme palette; can be overridden from CLI via --color.
        self.main_bg_color = RGBColor(141, 175, 208)
        self.main_inactive_text = RGBColor(245, 248, 252)
        self.main_active_chip_bg = RGBColor(189, 212, 234)
        self.main_active_text = RGBColor(43, 109, 180)
        self.sub_inactive_text = RGBColor(154, 154, 154)
        self.sub_active_text = RGBColor(43, 109, 180)
        self.sub_line_color = RGBColor(130, 175, 220)

        # Base font size in points for navigation and body
        self.font_size_pt = font_size if font_size and font_size > 0 else 22.0
        self.sub_font_size_pt = max(self.font_size_pt * 0.84, 12.0)
        # Scale navigation row heights with font size (22pt baseline).
        base_main_nav_height_in = 0.56
        base_sub_nav_height_in = 0.30
        scale = self.font_size_pt / 22.0
        self.main_nav_row_height = Inches(base_main_nav_height_in * scale)
        self.sub_nav_row_height = Inches(base_sub_nav_height_in * scale)
        self.sub_nav_side_margin = Inches(0.14)
        self.sub_nav_line_thickness = Inches(0.03)
        self.sub_nav_label_gap = Inches(0.1)
        # slide_width/slide_height from python-pptx are int-like EMU values.
        # Keep them as concrete ints to avoid Optional math issues in type checkers.
        self._slide_width: int = 0
        self._slide_height: int = 0

        # Default fonts for body textboxes: Latin and East Asian.
        # Note: PowerPoint uses separate font slots; setting only `font.name`
        # often doesn't affect CJK rendering, so we also set `a:ea` explicitly.
        self.body_font_latin = "Times New Roman"
        self.body_font_east_asian = "標楷體"

    def _set_paragraph_default_fonts(self, paragraph) -> None:
        pPr = paragraph._p.get_or_add_pPr()
        defRPr = pPr.get_or_add_defRPr()
        defRPr.get_or_add_latin().typeface = self.body_font_latin

        ea = defRPr.find(qn("a:ea"))
        if ea is None:
            ea = OxmlElement("a:ea")
            defRPr.append(ea)
        ea.set("typeface", self.body_font_east_asian)

    def build(
        self,
        outline: Outline,
        output_path: Path,
        template_path: Path,
    ) -> None:
        prs = PresentationFactory(str(template_path))
        # python-pptx stubs may type these as Optional/Unknown; guard for type checkers.
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        if slide_width is None or slide_height is None:
            raise ValueError("Presentation slide dimensions are not set.")
        self._slide_width = int(slide_width)
        self._slide_height = int(slide_height)
        for plan_entry in outline.iter_slide_plan():
            self._add_slide(prs, outline.sections, plan_entry)
        prs.save(str(output_path))

    def _add_slide(
        self,
        prs: PptxPresentation,
        sections: Iterable[OutlineItem],
        plan_entry: SlidePlanEntry,
    ) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        nav_bottom = self._add_navigation(slide, sections, plan_entry)
        self._add_body_placeholder(slide, plan_entry, nav_bottom)

    def _add_navigation(
        self,
        slide,
        sections: Iterable[OutlineItem],
        plan_entry: SlidePlanEntry,
    ) -> int:
        current_section = plan_entry.section
        current_child = plan_entry.child
        top = int(self.nav_top_margin)
        section_titles = [section.title for section in sections]
        top = self._draw_main_navigation_row(slide, section_titles, current_section.title, top)
        if current_section.children:
            child_titles = [child.title for child in current_section.children]
            active_child = current_child.title if current_child else None
            top = self._draw_sub_navigation_row(slide, child_titles, active_child, top)
        return top

    def _draw_main_navigation_row(
        self, slide, titles, active_title: Optional[str], top: int
    ) -> int:
        if not titles:
            return top
        if self._slide_width <= 0:
            raise ValueError("Slide width is not set.")
        row_height = int(self.main_nav_row_height)
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(self.nav_side_margin),
            top,
            int(self._slide_width - int(self.nav_side_margin) * 2),
            row_height,
        )
        self._style_solid_shape(bg, self.main_bg_color)

        count = len(titles)
        base_tab_width = int(self._slide_width // count)
        for idx, title in enumerate(titles):
            left = idx * base_tab_width
            if idx == count - 1:
                item_width = int(self._slide_width - left)
            else:
                item_width = base_tab_width

            if title == active_title:
                # Make the chip taller and size width by title length while
                # keeping it inside the tab area.
                chip_margin_y = max(int(row_height * 0.12), int(Inches(0.02)))
                chip_height = int(row_height - chip_margin_y * 2)
                chip_padding_x = int(Inches(0.24))
                estimated_text_width = self._estimate_text_width_emu(title, self.font_size_pt)
                desired_chip_width = int(estimated_text_width + chip_padding_x * 2)
                max_chip_width = int(max(item_width - int(Inches(0.1)), int(item_width * 0.6)))
                min_chip_width = int(min(item_width, max(estimated_text_width, int(Inches(0.95)))))
                chip_width = max(min(desired_chip_width, max_chip_width), min_chip_width)
                chip_margin_x = int(max((item_width - chip_width) // 2, 0))
                chip = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    int(left + chip_margin_x),
                    int(top + chip_margin_y),
                    chip_width,
                    chip_height,
                )
                self._style_solid_shape(chip, self.main_active_chip_bg)
                chip.adjustments[0] = 0.2

            self._add_centered_label(
                slide,
                title,
                int(left),
                top,
                int(item_width),
                row_height,
                self.main_active_text if title == active_title else self.main_inactive_text,
            )
        return top + row_height

    def _draw_sub_navigation_row(self, slide, titles, active_title: Optional[str], top: int) -> int:
        if self._slide_width <= 0:
            raise ValueError("Slide width is not set.")

        row_height = int(self.sub_nav_row_height)
        line_thickness = int(self.sub_nav_line_thickness)
        center_y = int(top + row_height // 2)
        side_margin = int(self.sub_nav_side_margin)

        left_line_start = side_margin
        left_line_end = int(left_line_start + int(Inches(0.40)))
        left_line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left_line_start,
            int(center_y),
            int(left_line_end - left_line_start),
            line_thickness,
        )
        self._style_solid_shape(left_line, self.sub_line_color)

        labels_left = int(left_line_end + int(self.sub_nav_label_gap))
        labels_right = int(self._slide_width - side_margin - int(Inches(0.35)))
        cursor = labels_left
        item_gap = int(Inches(0.08))
        min_label_width = int(Inches(0.35))

        if titles and labels_right > labels_left:
            text_padding = int(Inches(0.08))
            preferred_widths = [
                max(
                    self._estimate_text_width_emu(title, self.sub_font_size_pt) + text_padding,
                    min_label_width,
                )
                for title in titles
            ]
            target_widths = self._fit_widths_to_space(
                preferred_widths,
                int(labels_right - labels_left),
                item_gap,
                min_label_width,
            )
            for idx, title in enumerate(titles):
                remaining = max(labels_right - cursor, 0)
                if remaining <= 0:
                    break
                width = min(target_widths[idx], remaining)
                self._add_left_label(
                    slide,
                    title,
                    cursor,
                    top,
                    width,
                    row_height,
                    self.sub_active_text if title == active_title else self.sub_inactive_text,
                    self.sub_font_size_pt,
                )
                cursor += int(width + item_gap)

        right_line_start = int(min(max(cursor, labels_left), self._slide_width - side_margin))
        right_line_end = int(self._slide_width - side_margin)
        if right_line_end > right_line_start:
            right_line = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                right_line_start,
                int(center_y),
                int(right_line_end - right_line_start),
                line_thickness,
            )
            self._style_solid_shape(right_line, self.sub_line_color)

        return top + row_height

    def _style_solid_shape(self, shape, fill_color: RGBColor) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        shape.shadow.inherit = False

    def _add_centered_label(
        self,
        slide,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        color: RGBColor,
    ) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(self.font_size_pt)
        run.font.color.rgb = color
        run.font.name = self.body_font_latin

    def _add_left_label(
        self,
        slide,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        color: RGBColor,
        font_size_pt: float,
    ) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = color
        run.font.name = self.body_font_latin

    def _estimate_text_width_emu(self, text: str, font_size_pt: float) -> int:
        # Estimate width by character category to reduce over/under-estimation
        # for mixed-case English, digits, spaces, and CJK text.
        stripped = text.strip()
        if not stripped:
            stripped = " "

        em_width = 0.0
        for ch in stripped:
            if ch.isspace():
                em_width += 0.32
            elif self._is_cjk_char(ch):
                em_width += 1.0
            elif ch in "ilI.,:;'`!|":
                em_width += 0.3
            elif ch in "MW@#%&":
                em_width += 0.8
            elif ch.isupper():
                em_width += 0.65
            elif ch.islower():
                em_width += 0.5
            elif ch.isdigit():
                em_width += 0.55
            else:
                em_width += 0.5

        width_pt = max(em_width * font_size_pt, font_size_pt * 1.2)
        return int(width_pt * 12700)

    def _is_cjk_char(self, ch: str) -> bool:
        return unicodedata.east_asian_width(ch) in {"W", "F"}

    def _fit_widths_to_space(
        self,
        preferred_widths: list[int],
        available_width: int,
        item_gap: int,
        min_width: int,
    ) -> list[int]:
        count = len(preferred_widths)
        if count == 0:
            return []

        usable = max(available_width - item_gap * max(count - 1, 0), count)
        preferred_total = sum(preferred_widths)
        if preferred_total <= usable:
            return preferred_widths

        scale = usable / preferred_total
        scaled = [max(min_width, int(round(width * scale))) for width in preferred_widths]

        total = sum(scaled)
        while total > usable:
            idx = max(range(count), key=lambda i: scaled[i])
            if scaled[idx] <= min_width:
                break
            scaled[idx] -= 1
            total -= 1

        return scaled

    def _add_body_placeholder(
        self,
        slide,
        plan_entry: SlidePlanEntry,
        nav_bottom: int,
    ) -> None:
        if self._slide_width <= 0 or self._slide_height <= 0:
            raise ValueError("Slide dimensions are not set.")
        body_top = int(nav_bottom + int(self.body_margin_top))
        height = int(self._slide_height - body_top - int(Inches(0.5)))
        width = int(self._slide_width - int(self.body_side_margin) * 2)
        box = slide.shapes.add_textbox(int(self.body_side_margin), body_top, width, height)
        tf = box.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        tf.paragraphs[0].space_after = 0
        tf.paragraphs[0].space_before = 0
        tf.paragraphs[0].line_spacing = 1.1
        para.font.size = Pt(self.font_size_pt)
        self._set_paragraph_default_fonts(para)
        para.text = ""
        extra_para = tf.add_paragraph()
        extra_para.font.size = Pt(self.font_size_pt)
        self._set_paragraph_default_fonts(extra_para)
